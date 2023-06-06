##############################################################
# CODE WITH GOZ Youtube Channel - Python Video No.4 - Script 1
##############################################################
# Description: Simple script for diploma generation
# Dependencies:
# ----------------------------------------
# pip install python-pptx
# ----------------------------------------
# Author: Goz
# https://codewithgoz.com
##############################################################

import csv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE


# Define variables for diploma
img_path = "bg.png"
message1 = "Agradece a:"
message2 = "Por su participación en el Concurso:"
event_name = "Come Tacos Hasta Reventar, Edición 2023"
event_date = "Realizado el 1 y 2 de Junio de 2023"
event_place = "En las instalaciones de Taquitos Don Gustavo"
expedition_date = "Ciudad de México, 3 de Junio de 2023"
responsible = "Gustavo Gómez Macías (Goz)"
responsible_title = "Dueño de la Sucursal"
participants_path = "names.csv"
my_font = 'Abel'

def add_my_textbox(slide, message, font_size, my_font, R, G, B, left, top, width, height):
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = text_box.text_frame.paragraphs[0]
    tf.alignment = PP_ALIGN.CENTER
    tf.text = message
    font = tf.font
    font.name = my_font
    font.size = Pt(font_size)
    font.color.rgb = RGBColor(R,G,B)

# Create presentation
my_presentation = Presentation()

# Open participants csv file
participants = csv.reader(open(participants_path))

# Loop through participants generating one diploma per participant
for participant in participants:
    blank_slide_layout = my_presentation.slide_layouts[6]
    slide = my_presentation.slides.add_slide(blank_slide_layout)
    bg = slide.shapes.add_picture(img_path, Inches(0), Inches(-0.1), width = Inches(10))
    add_my_textbox(slide, message1, 24, my_font, 0x54, 0x54, 0x53, 0, 2.60, 10, 2)
    add_my_textbox(slide, participant[0], 32, my_font, 0xe2, 0xa1, 0x4f, 0, 3.26, 10, 2)
    add_my_textbox(slide, message2, 16, my_font, 0x32, 0x33, 0x32, 0, 4.16, 10, 2)
    add_my_textbox(slide, event_name, 20, my_font, 0xe2, 0xa1, 0x4f, 0, 4.65, 10, 2)
    add_my_textbox(slide, event_date, 16, my_font, 0x54, 0x54, 0x53, 0, 5.2, 10, 2)
    add_my_textbox(slide, event_place, 16, my_font, 0x54, 0x54, 0x53, 0, 5.48, 10, 2)
    add_my_textbox(slide, expedition_date, 12, my_font, 0x32, 0x33, 0x32, 0, 6.0, 10, 2)
    add_my_textbox(slide, responsible, 12, my_font, 0x32, 0x33, 0x32, 0, 6.8, 10, 2)
    add_my_textbox(slide, responsible_title, 12, my_font, 0x32, 0x33, 0x32, 0, 7, 10, 2)
    shape = slide.shapes.add_shape( MSO_LINE.SOLID, Inches(3.30), Inches(6.78), Inches(3.5), Inches(0.01))
    line = shape.line
    line.color.rgb = RGBColor(0,0,0)

# Save presentation
my_presentation.save('concursotacos.pptx')




















